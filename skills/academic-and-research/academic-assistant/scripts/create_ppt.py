from pptx import Presentation
from pptx.util import Inches, Pt
import sys
import json

def create_presentation(data, output_file="report.pptx"):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = data.get("title", "學術研究報告")
    subtitle.text = data.get("subtitle", "由學術研究助手自動生成")
    
    # Content Slides
    for item in data.get("slides", []):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = item.get("header", "")
        tf = body_shape.text_frame
        tf.text = item.get("content", "")
        
    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 create_ppt.py <json_data_file>")
        sys.exit(1)
    
    with open(sys.argv[1], 'r') as f:
        data = json.load(f)
    
    create_presentation(data)
